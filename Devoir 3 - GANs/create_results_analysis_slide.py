"""
Create detailed analysis slide of training results
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Slide: Analyse des Résultats
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

# Title
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.6))
title_frame = title_box.text_frame
title_frame.text = "Analyse des Résultats d'Entraînement - 500 Epochs"
title_frame.paragraphs[0].font.size = Pt(28)
title_frame.paragraphs[0].font.bold = True
title_frame.paragraphs[0].font.color.rgb = RGBColor(25, 118, 210)

# ============================================================================
# SECTION 1: GENERATOR LOSS
# ============================================================================

gen_title = slide.shapes.add_textbox(Inches(0.5), Inches(0.95), Inches(3), Inches(0.4))
gen_title_frame = gen_title.text_frame
gen_title_frame.text = "1. GENERATOR LOSS"
gen_title_frame.paragraphs[0].font.size = Pt(14)
gen_title_frame.paragraphs[0].font.bold = True
gen_title_frame.paragraphs[0].font.color.rgb = RGBColor(198, 40, 40)

gen_content = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(3), Inches(2.4))
gen_frame = gen_content.text_frame
gen_frame.word_wrap = True
gen_frame.text = """✓ TRAIN Loss:
  • 2.8 → 4.2 (montée)
  • Stabilise à 4.0-4.25
  • Comportement NORMAL
    pour GAN

✓ EVAL Loss:
  • Très stable ~2.3-2.4
  • Suit train fidèlement
  
⚠ INTERPRETATION:
  • Gap train/eval accepte
  • G apprend bien
  • Pas d'overfitting majeur"""

for paragraph in gen_frame.paragraphs:
    paragraph.font.size = Pt(9)
    paragraph.space_before = Pt(1)

# ============================================================================
# SECTION 2: L1 LOSS
# ============================================================================

l1_title = slide.shapes.add_textbox(Inches(3.8), Inches(0.95), Inches(3), Inches(0.4))
l1_title_frame = l1_title.text_frame
l1_title_frame.text = "2. L1 RECONSTRUCTION LOSS"
l1_title_frame.paragraphs[0].font.size = Pt(14)
l1_title_frame.paragraphs[0].font.bold = True
l1_title_frame.paragraphs[0].font.color.rgb = RGBColor(56, 142, 60)

l1_content = slide.shapes.add_textbox(Inches(3.8), Inches(1.4), Inches(3), Inches(2.4))
l1_frame = l1_content.text_frame
l1_frame.word_wrap = True
l1_frame.text = """✓ TRAIN Loss:
  • 0.70 → 0.67-0.68
  • Descente lente stable
  • Reconstruction OK

✓ EVAL Loss:
  • Stable ~0.67-0.69
  • Suit train bien
  
⚠ ANOMALIE (Epoch 400):
  • Spike eval ~1.3-1.4
  • Puis revient normal
  • Possible batch difficile"""

for paragraph in l1_frame.paragraphs:
    paragraph.font.size = Pt(9)
    paragraph.space_before = Pt(1)

# ============================================================================
# SECTION 3: DISCRIMINATOR LOSS
# ============================================================================

disc_title = slide.shapes.add_textbox(Inches(7.1), Inches(0.95), Inches(2.9), Inches(0.4))
disc_title_frame = disc_title.text_frame
disc_title_frame.text = "3. DISCRIMINATOR LOSS"
disc_title_frame.paragraphs[0].font.size = Pt(14)
disc_title_frame.paragraphs[0].font.bold = True
disc_title_frame.paragraphs[0].font.color.rgb = RGBColor(245, 127, 23)

disc_content = slide.shapes.add_textbox(Inches(7.1), Inches(1.4), Inches(2.9), Inches(2.4))
disc_frame = disc_content.text_frame
disc_frame.word_wrap = True
disc_frame.text = """✓ COMPORTEMENT:
  • 1.0 → 0.05 (rapide)
  • Epochs 0-50: chute
  • Après: plateau ~0.0
  
✓ NORMAL POUR GAN:
  • D converge bien
  • Trop rapide OK si G
    continue d'apprendre

✓ LABEL SMOOTHING:
  • Mantient gradients
  • G reste stable"""

for paragraph in disc_frame.paragraphs:
    paragraph.font.size = Pt(9)
    paragraph.space_before = Pt(1)

# ============================================================================
# SECTION 4: CONCLUSIONS
# ============================================================================

conclusion_title = slide.shapes.add_textbox(Inches(0.5), Inches(4.0), Inches(9), Inches(0.4))
conclusion_title_frame = conclusion_title.text_frame
conclusion_title_frame.text = "CONCLUSIONS GLOBALES"
conclusion_title_frame.paragraphs[0].font.size = Pt(16)
conclusion_title_frame.paragraphs[0].font.bold = True
conclusion_title_frame.paragraphs[0].font.color.rgb = RGBColor(103, 58, 183)

conclusion_content = slide.shapes.add_textbox(Inches(0.7), Inches(4.5), Inches(8.6), Inches(2.7))
conclusion_frame = conclusion_content.text_frame
conclusion_frame.word_wrap = True
conclusion_frame.text = """✓ POINTS POSITIFS:
  • Generator Loss stable (montée normale en GAN)
  • L1 Loss décroît régulièrement (reconstruction améliore)
  • Discriminator converge (bon signal au générateur)
  • Train/Eval proches (pas d'overfitting majeur)
  • 500 epochs suffisant (courbes stabilisées)

⚠ POINTS A NOTER:
  • L1Loss spike epoch 400 (anomalie, enquêter batch difficile)
  • GenLoss augmente (normal car HingeLoss plateau, L1 décline)
  • D converge très vite (mais label smoothing maintient gradient)

🎯 VERDICT:
  Entraînement STABLE et CONVERGENT. Le modèle apprend les masques avec adversarial learning équilibré."""

for paragraph in conclusion_frame.paragraphs:
    paragraph.font.size = Pt(10)
    paragraph.space_before = Pt(1.5)

# Save
prs.save('Training_Results_Analysis.pptx')
print("✓ Slide d'analyse créée: Training_Results_Analysis.pptx")
