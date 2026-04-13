"""
Create a detailed PowerPoint slide explaining Label Smoothing
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Slide: Label Smoothing Details
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

# Title
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
title_frame = title_box.text_frame
title_frame.text = "Label Smoothing (α = 0.2)"
title_frame.paragraphs[0].font.size = Pt(40)
title_frame.paragraphs[0].font.bold = True
title_frame.paragraphs[0].font.color.rgb = RGBColor(25, 118, 210)

# ============================================================================
# POINT 1: Évite l'Overfitting du Discriminateur
# ============================================================================

point1_title = slide.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(8.4), Inches(0.5))
point1_title_frame = point1_title.text_frame
point1_title_frame.text = "• Point 1: Évite l'Overfitting du Discriminateur"
point1_title_frame.paragraphs[0].font.size = Pt(18)
point1_title_frame.paragraphs[0].font.bold = True
point1_title_frame.paragraphs[0].font.color.rgb = RGBColor(198, 40, 40)

point1_content = slide.shapes.add_textbox(Inches(1.2), Inches(1.9), Inches(8), Inches(2.2))
point1_frame = point1_content.text_frame
point1_frame.word_wrap = True
point1_frame.text = """SANS Label Smoothing (α = 0):
  • Discriminateur cible: D(real) = 1.0 exactement
  • D devient extremement confiant (overfit)
  • Loss = max(0, 1 - D(real)) avec D(real) >> 1
  • D converge trop vite (plateau d'apprentissage)

AVEC Label Smoothing (α = 0.2):
  • Discriminateur cible: D(real) = 0.8 (au lieu de 1.0)
  • Loss = max(0, 0.8 - D(real))
  • D continue à apprendre meme avec des valeurs hautes
  • Reste un defi modere, pas un plateau"""

for paragraph in point1_frame.paragraphs:
    paragraph.font.size = Pt(10)
    paragraph.space_before = Pt(2)

# ============================================================================
# POINT 2: Maintient le Gradient pour le Générateur
# ============================================================================

point2_title = slide.shapes.add_textbox(Inches(0.8), Inches(4.2), Inches(8.4), Inches(0.5))
point2_title_frame = point2_title.text_frame
point2_title_frame.text = "• Point 2: Maintient un Gradient Stable pour le Générateur"
point2_title_frame.paragraphs[0].font.size = Pt(18)
point2_title_frame.paragraphs[0].font.bold = True
point2_title_frame.paragraphs[0].font.color.rgb = RGBColor(56, 142, 60)

point2_content = slide.shapes.add_textbox(Inches(1.2), Inches(4.8), Inches(8), Inches(2.4))
point2_frame = point2_content.text_frame
point2_frame.word_wrap = True
point2_frame.text = """PROBLEME SANS LABEL SMOOTHING:
  • D(real) atteint 5.0 ou 10.0 (D devient trop puissant)
  • Loss Hinge = max(0, 1 - 10) = 0 (aucune penalite!)
  • Gradient est 0 → D arrete d'apprendre
  • G reçoit moins de feedback du discriminateur (gradient faible)
  • G stagne, ne produit plus d'ameliorations

SOLUTION AVEC LABEL SMOOTHING:
  • Cible est D(real) = 0.8 (moins extreme)
  • Meme si D(real) = 2.0, loss = max(0, 0.8 - 2.0) = 0
  • MAIS gradient reste constant (Hinge loss propriete)
  • D continue a donner du feedback utile a G
  • G continue d'apprendre et d'ameliorer les masques"""

for paragraph in point2_frame.paragraphs:
    paragraph.font.size = Pt(10)
    paragraph.space_before = Pt(2)

# Save
prs.save('Label_Smoothing_Detailed.pptx')
print("✓ Slide détaillée créée: Label_Smoothing_Detailed.pptx")
