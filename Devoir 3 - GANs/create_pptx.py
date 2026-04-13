"""
Script pour crear un PowerPoint de présentation Pix2Pix
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

###############################################################################
# Slide 1: Title Slide
###############################################################################
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
background = slide1.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(25, 118, 210)  # Blue

title_box = slide1.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.5))
title_frame = title_box.text_frame
title_frame.text = "Pix2Pix GAN"
title_frame.paragraphs[0].font.size = Pt(60)
title_frame.paragraphs[0].font.bold = True
title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

subtitle_box = slide1.shapes.add_textbox(Inches(1), Inches(4), Inches(8), Inches(2))
subtitle_frame = subtitle_box.text_frame
subtitle_frame.text = "Détection de Défauts de Câbles\nAvec Réseaux antagonistes conditionnels"
subtitle_frame.paragraphs[0].font.size = Pt(24)
subtitle_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

###############################################################################
# Slide 2: Architecture - Complete Overview
###############################################################################
slide2 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank
title2 = slide2.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
title2.text_frame.text = "Architecture: U-Net Generator + PatchGAN Discriminator"
title2.text_frame.paragraphs[0].font.size = Pt(32)
title2.text_frame.paragraphs[0].font.bold = True

# Add image (pix2pix_training_losses.png if available)
try:
    slide2.shapes.add_picture('pix2pix_architecture_diagram.png', Inches(0.5), Inches(1), width=Inches(9))
except:
    # Fallback: Add text description
    desc_box = slide2.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5))
    desc_frame = desc_box.text_frame
    desc_frame.word_wrap = True
    desc_frame.text = """• GENERATOR (U-Net):
  - Encodeur: Downsampling 128 → 64 → 32 → 16 → 8
  - Bottleneck: 8×8 feature maps
  - Décodeur: Upsampling avec skip connections
  - Output: Sigmoid [0,1] pour masques binaires

• DISCRIMINATOR (PatchGAN):
  - Classifie 70×70 patches comme réels/faux
  - 4 bloc de downsampling avec dropout
  - Output: 7×7 heatmap de logits
  
• LOSS COMBINÉE:
  - Hinge Loss (adversarial)
  - BCE + L1 (reconstruction)"""
    for paragraph in desc_frame.paragraphs:
        paragraph.font.size = Pt(16)

###############################################################################
# Slide 3: Loss Functions
###############################################################################
slide3 = prs.slides.add_slide(prs.slide_layouts[5])
title3 = slide3.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
title3.text_frame.text = "Loss Functions Utilisées"
title3.text_frame.paragraphs[0].font.size = Pt(32)
title3.text_frame.paragraphs[0].font.bold = True

loss_box = slide3.shapes.add_textbox(Inches(1), Inches(1.2), Inches(8), Inches(5.8))
loss_frame = loss_box.text_frame
loss_frame.word_wrap = True
loss_frame.text = """DISCRIMINATOR LOSS (Hinge Loss):
  L_D = max(0, 1 - D(real)) + max(0, 1 + D(fake))
  • Pénalise quand D ne discrimine pas assez
  • Plus stable que BCE Loss

GENERATOR LOSS:
  L_G = max(0, 1 - D(fake)) + λ_L1 × (BCE + 0.5×L1)
  
  Composants:
  • Hinge Loss: Force le générateur à tromper le discriminateur
  • BCE Loss: Classification binaire des masques
  • L1 Loss: Régularisation spatiale (smoothness)
  
HYPERPARAMÈTRES:
  • λ_L1 = 2.0 (poids de la reconstruction L1)
  • label_smooth = 0.2 (smoothing du discriminateur)"""

for paragraph in loss_frame.paragraphs:
    paragraph.font.size = Pt(13)

###############################################################################
# Slide 4: Problèmes Rencontrés et Solutions
###############################################################################
slide4 = prs.slides.add_slide(prs.slide_layouts[5])
title4 = slide4.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
title4.text_frame.text = "Problèmes Rencontrés et Solutions"
title4.text_frame.paragraphs[0].font.size = Pt(32)
title4.text_frame.paragraphs[0].font.bold = True

problems_box = slide4.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5.8))
problems_frame = problems_box.text_frame
problems_frame.word_wrap = True
problems_frame.text = """1. MODE COLLAPSE (Outputs Gris)
   Problème: Générateur produit toujours des masques gris (0.3-0.4)
   Causes: 
   • L1 Loss force mathématiquement la moyenne pour cibles binaires
   • Tanh activation [-1,1] inadapté pour masques [0,1]
   Solutions appliquées:
   • Sigmoid output [0,1] au lieu de Tanh
   • Réduit λ_L1 pour réduire pression L1
   • Augmenté learning rate discriminateur (stabilité)

2. DISCRIMINATEUR CONVERGE TROP VITE
   Problème: DiscLoss → 0 rapidement, perte de signal gradient
   Solutions:
   • Hinge Loss au lieu de BCE (gradient plus stable)
   • Label smoothing (real_labels = 0.8 au lieu de 1.0)
   • Dropout 0.3 sur discriminateur (régularisation)

3. ENTRAÎNEMENT INSTABLE
   Problème: GenLoss augmente mais L1Loss diminue légèrement
   Solution: Plus d'epochs (200) pour convergence plus lente"""

for paragraph in problems_frame.paragraphs:
    paragraph.font.size = Pt(11)

###############################################################################
# Slide 5: Résultats - Loss Curves
###############################################################################
slide5 = prs.slides.add_slide(prs.slide_layouts[5])
title5 = slide5.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
title5.text_frame.text = "Résultats - Courbes de Loss"
title5.text_frame.paragraphs[0].font.size = Pt(32)
title5.text_frame.paragraphs[0].font.bold = True

# Try to add loss plot
try:
    slide5.shapes.add_picture('pix2pix_training_losses.png', Inches(0.8), Inches(1.1), width=Inches(8.4))
except FileNotFoundError:
    info_box = slide5.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5))
    info_frame = info_box.text_frame
    info_frame.word_wrap = True
    info_frame.text = """Les courbes de loss montrent:
    
    • GenLoss: Stable autour de 2.75-4.5 (comportement normal en GAN)
    • DiscLoss: Converge rapidement (bon discriminateur)
    • L1Loss: Décroît progressivement (reconstruction améliore)
    
    Train vs Eval:
    • Train et Eval sont proches → pas d'overfitting majeur
    • Eval L1Loss stable → généralization acceptable"""
    for paragraph in info_frame.paragraphs:
        paragraph.font.size = Pt(16)

###############################################################################
# Slide 6: Résultats - Exemples de Génération
###############################################################################
slide6 = prs.slides.add_slide(prs.slide_layouts[5])
title6 = slide6.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
title6.text_frame.text = "Résultats - Exemples Générés vs Ground Truth"
title6.text_frame.paragraphs[0].font.size = Pt(32)
title6.text_frame.paragraphs[0].font.bold = True

# Try to add samples
try:
    slide6.shapes.add_picture('pix2pix_samples.png', Inches(0.8), Inches(1.1), width=Inches(8.4))
except FileNotFoundError:
    info_box2 = slide6.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5))
    info_frame2 = info_box2.text_frame
    info_frame2.word_wrap = True
    info_frame2.text = """Observations des résultats:

    ✓ POINTS POSITIFS:
      • Masques générés proches spatialement du ground truth
      • Détecte la forme générale des défauts
      • Outputs condition-dépendants (différentes entrées → différentes sorties)
      
    ⚠ AMÉLIORATIONS NÉCESSAIRES:
      • Output range [0.2-0.5] au lieu de [0,1] (trop gris)
      • Manque de netteté/contraste dans les masques
      • Qualité variable selon type de défaut"""
    for paragraph in info_frame2.paragraphs:
        paragraph.font.size = Pt(14)

###############################################################################
# Slide 7: Test de Robustesse au Bruit
###############################################################################
slide7 = prs.slides.add_slide(prs.slide_layouts[5])
title7 = slide7.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
title7.text_frame.text = "Robustesse au Bruit - Condition Dependence"
title7.text_frame.paragraphs[0].font.size = Pt(32)
title7.text_frame.paragraphs[0].font.bold = True

robustness_box = slide7.shapes.add_textbox(Inches(1), Inches(1.2), Inches(8), Inches(5.8))
robustness_frame = robustness_box.text_frame
robustness_frame.word_wrap = True
robustness_frame.text = """TEST: Robustesse au Bruit Gaussien

RÉSULTATS:
  • Entrée réelle (image test) → Output structuré
  • Bruit aléatoire (Gaussian noise) → Output dégradé (1.5-2x pire L1 error)
  • Bruit progressif → Dégradation progressive

CONCLUSION:
  ✓ Le modèle apprend bien la dépendance condition → sortie
  ✓ N'est pas une simple mémorisation (ne génère pas du bruit aléatoire)
  ✓ Utilise les features de l'image d'entrée pour générer les masques
  
IMPLICATION ACADÉMIQUE:
  Le générateur a appris des feature maps significatives liées à la 
  détection de défauts, plutôt que d'apprendre une sortie fixe."""

for paragraph in robustness_frame.paragraphs:
    paragraph.font.size = Pt(12)

###############################################################################
# Save presentation
###############################################################################
prs.save('Pix2Pix_Presentation.pptx')
print("✓ Présentation créée avec succès: Pix2Pix_Presentation.pptx")
